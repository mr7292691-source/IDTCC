#!/usr/bin/env python3
"""Generate the IDTCC hackathon deck (12 slides) → IDTCC_Hackathon.pptx.

Self-contained: only requires python-pptx. Run:

    python scripts/generate_ppt.py [--out IDTCC_Hackathon.pptx] [--team "Name"]

Edit TEAM / CONTACT / links below before presenting.
"""
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Brand palette ─────────────────────────────────────────────────────────────
BG      = RGBColor(0x0B, 0x0E, 0x14)   # near-black navy
PANEL   = RGBColor(0x14, 0x1A, 0x24)   # panel
ACCENT  = RGBColor(0xE8, 0x11, 0x2D)   # AMD/BMW-M red
ACCENT2 = RGBColor(0x2E, 0x9B, 0xF0)   # blue
GOLD    = RGBColor(0xF5, 0xB5, 0x0C)
WHITE   = RGBColor(0xF2, 0xF4, 0xF8)
MUTED   = RGBColor(0x9A, 0xA4, 0xB2)
GREEN   = RGBColor(0x2E, 0xCB, 0x70)

FONT = "Segoe UI"
W, H = Inches(13.333), Inches(7.5)


def _slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def _box(s, x, y, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line; shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def _text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=4):
    """runs: list of (text, size, color, bold) — each becomes a paragraph."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = FONT
    return tb


def _accent_bar(s, x, y, w=Inches(0.9), h=Pt(5)):
    _box(s, x, y, w, h, fill=ACCENT)


def _kicker(s, text):
    _text(s, Inches(0.7), Inches(0.45), Inches(11), Inches(0.4),
          [(text, 14, ACCENT, True)])


def _title(s, text, size=34):
    _text(s, Inches(0.7), Inches(0.78), Inches(12), Inches(1.0),
          [(text, size, WHITE, True)])
    _accent_bar(s, Inches(0.72), Inches(1.55))


def _footer(s, n):
    _text(s, Inches(0.7), Inches(7.02), Inches(8), Inches(0.4),
          [("IDTCC — Insurance Digital Twin Command Center", 9, MUTED, False)])
    _text(s, Inches(12.0), Inches(7.02), Inches(1.0), Inches(0.4),
          [(f"{n:02d} / 12", 9, MUTED, False)], align=PP_ALIGN.RIGHT)


def _bullets(s, x, y, w, items, size=15, gap=8, color=WHITE, marker="—"):
    runs = [(f"{marker}  {it}", size, color, False) for it in items]
    _text(s, x, y, w, Inches(4.5), runs, space_after=gap)


def _chip(s, x, y, label, w=Inches(2.2), h=Inches(0.55), fill=PANEL,
          fg=WHITE, line=ACCENT2):
    _box(s, x, y, w, h, fill=fill, line=line, line_w=1.25, round_=True)
    _text(s, x, y, w, h, [(label, 12, fg, True)],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _metric(s, x, y, value, label, w=Inches(2.7), color=ACCENT2):
    _box(s, x, y, w, Inches(1.5), fill=PANEL, line=color, line_w=1.25, round_=True)
    _text(s, x, y + Inches(0.18), w, Inches(0.8), [(value, 30, color, True)],
          align=PP_ALIGN.CENTER)
    _text(s, x, y + Inches(0.98), w, Inches(0.45), [(label, 11, MUTED, False)],
          align=PP_ALIGN.CENTER)


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_title(prs, team, contact):
    s = _slide(prs)
    _box(s, 0, Inches(2.55), W, Pt(6), fill=ACCENT)
    _text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(1.0),
          [("AMD AI HACKATHON", 16, ACCENT, True)])
    _text(s, Inches(0.7), Inches(2.75), Inches(12.2), Inches(1.6),
          [("IDTCC", 72, WHITE, True)])
    _text(s, Inches(0.72), Inches(4.0), Inches(12), Inches(0.8),
          [("Insurance Digital Twin Command Center", 26, WHITE, False)])
    _text(s, Inches(0.72), Inches(4.8), Inches(12), Inches(0.6),
          [("AI catastrophe simulation & insurance risk intelligence — on AMD MI300X",
            15, MUTED, False)])
    _chip(s, Inches(0.72), Inches(5.7), "50,000 Digital Twins", w=Inches(2.7))
    _chip(s, Inches(3.6), Inches(5.7), "8 AI Agents", w=Inches(2.0), line=GOLD)
    _chip(s, Inches(5.8), Inches(5.7), "35 Cities · 14 States", w=Inches(2.8), line=GREEN)
    _text(s, Inches(0.72), Inches(6.5), Inches(12), Inches(0.5),
          [(f"Team {team}   ·   {contact}", 12, MUTED, False)])


def slide_challenge(prs):
    s = _slide(prs)
    _kicker(s, "INDUSTRY CHALLENGE")
    _title(s, "Insurers react after the storm — not before")
    _bullets(s, Inches(0.7), Inches(1.95), Inches(6.0), [
        "Indian cyclones drive billions in insured catastrophe losses every year.",
        "Loss estimation today is top-down, slow, and post-landfall.",
        "Adjusters are deployed reactively → slow claims, higher leakage.",
        "Fraud spikes during catastrophes; reserves are set with limited granularity.",
        "Regulators demand explainable, auditable risk decisions.",
    ], size=15, gap=12)
    _box(s, Inches(7.1), Inches(1.95), Inches(5.5), Inches(4.4),
         fill=PANEL, line=ACCENT, line_w=1.25, round_=True)
    _text(s, Inches(7.4), Inches(2.2), Inches(5.0), Inches(0.6),
          [("Why Digital Twins", 18, ACCENT, True)])
    _bullets(s, Inches(7.4), Inches(2.9), Inches(5.0), [
        "Per-property risk: construction, flood zone, elevation.",
        "Social vulnerability: infants, elderly, disabled.",
        "Prior-claim & fraud history per twin.",
        "Bottom-up aggregation, not top-down estimates.",
        "Privacy-safe synthetic data → scales to 50K.",
    ], size=14, gap=11, color=WHITE)


def slide_architecture(prs):
    s = _slide(prs)
    _kicker(s, "SOLUTION ARCHITECTURE")
    _title(s, "Three services, one stateless pipeline")
    y = Inches(2.2)
    cols = [
        ("React 18 SPA", ":5173", "Recharts · Leaflet\nretry + SSE client", ACCENT2),
        ("FastAPI", ":8000", "REST + SSE\nLangGraph orchestration", ACCENT),
        ("vLLM", ":8001", "Qwen3-14B · bf16\nROCm", GOLD),
        ("AMD MI300X", "192 GB", "HBM3 · ROCm 6.x\nKV-cache bound", GREEN),
    ]
    x = Inches(0.7)
    cw = Inches(2.85)
    for name, port, sub, col in cols:
        _box(s, x, y, cw, Inches(1.9), fill=PANEL, line=col, line_w=1.5, round_=True)
        _text(s, x, y + Inches(0.2), cw, Inches(0.5), [(name, 18, col, True)],
              align=PP_ALIGN.CENTER)
        _text(s, x, y + Inches(0.72), cw, Inches(0.4), [(port, 13, WHITE, True)],
              align=PP_ALIGN.CENTER)
        _text(s, x, y + Inches(1.12), cw, Inches(0.7), [(sub, 11, MUTED, False)],
              align=PP_ALIGN.CENTER)
        if x + cw < Inches(12):
            _text(s, x + cw - Inches(0.05), y + Inches(0.7), Inches(0.5), Inches(0.5),
                  [("→", 22, WHITE, True)], align=PP_ALIGN.CENTER)
        x = Emu(int(x) + int(cw) + int(Inches(0.08)))
    _box(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(1.9),
         fill=PANEL, line=ACCENT2, line_w=1.25, round_=True)
    _text(s, Inches(1.0), Inches(4.7), Inches(11.3), Inches(0.5),
          [("Cross-cutting: confidence · explainability · guardrails · "
            "Prometheus metrics · LangSmith tracing · structured logging", 14, WHITE, True)])
    _bullets(s, Inches(1.0), Inches(5.35), Inches(11.0), [
        "Stateless backend → horizontal scaling (more uvicorn workers/instances). Deterministic, seeded twins.",
        "Resilience: tenacity retry · vLLM→Anthropic failover · frontend local fallback.",
        "Health/readiness endpoints + Prometheus metrics + CI with security scanning.",
    ], size=13, gap=7)


def slide_agents(prs):
    s = _slide(prs)
    _kicker(s, "AI AGENT ECOSYSTEM")
    _title(s, "8 agents orchestrated by a LangGraph DAG")
    agents = [
        ("1 · Weather", "severity + hazards"),
        ("2 · Risk Exposure", "portfolio in radius"),
        ("3 · Claims", "expected loss"),
        ("4 · Fraud (FAISS)", "anomaly detection"),
        ("5 · Reserve", "IBNR + cat buffer"),
        ("6 · Resource (K-Means)", "adjuster zones"),
        ("7 · Alerts", "personalised SMS"),
        ("8 · LLM-as-Judge", "audit + verdict"),
    ]
    x0, y0 = Inches(0.7), Inches(2.0)
    cw, ch = Inches(2.92), Inches(1.15)
    for i, (name, sub) in enumerate(agents):
        r, c = divmod(i, 4)
        x = Emu(int(x0) + c * (int(cw) + int(Inches(0.07))))
        y = Emu(int(y0) + r * (int(ch) + int(Inches(0.15))))
        col = GOLD if "Judge" in name else ACCENT2
        _box(s, x, y, cw, ch, fill=PANEL, line=col, line_w=1.25, round_=True)
        _text(s, x + Inches(0.15), y + Inches(0.16), cw - Inches(0.3), Inches(0.4),
              [(name, 14, col, True)])
        _text(s, x + Inches(0.15), y + Inches(0.62), cw - Inches(0.3), Inches(0.4),
              [(sub, 11, MUTED, False)])
    _box(s, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.2),
         fill=PANEL, line=ACCENT, line_w=1.25, round_=True)
    _text(s, Inches(1.0), Inches(5.38), Inches(11.3), Inches(1.0), [
        ("DAG:  Weather → Risk →  (Claims ∥ Fraud)  →  Reserve →  (Resource ∥ Alerts)  →  Judge → Forecast",
         15, WHITE, True),
        ("Parallel branches join at the Judge. Every agent returns confidence + explainability; "
         "a failed agent degrades gracefully — the run always completes.", 12, MUTED, False),
    ], space_after=8)


def slide_twins(prs):
    s = _slide(prs)
    _kicker(s, "DIGITAL TWIN ENGINE")
    _title(s, "50,000 living property twins")
    _metric(s, Inches(0.7),  Inches(2.1), "50,000", "property twins", color=ACCENT)
    _metric(s, Inches(3.55), Inches(2.1), "35", "cities", color=ACCENT2)
    _metric(s, Inches(6.4),  Inches(2.1), "14", "states", color=GOLD)
    _metric(s, Inches(9.25), Inches(2.1), "100+", "attributes / twin", color=GREEN)
    _text(s, Inches(0.7), Inches(3.9), Inches(6.0), Inches(0.5),
          [("Per-twin attributes", 17, ACCENT2, True)])
    _bullets(s, Inches(0.7), Inches(4.4), Inches(6.0), [
        "Geo: lat/lng, area, address, flood zone (A/B/C).",
        "Structure: construction & roof type, floors, year built, elevation.",
        "Exposure: sum insured, vulnerability index, claim probability.",
        "Social: infants / elderly / disabled flags.",
        "History: prior claims, prior fraud flags.",
    ], size=13, gap=9)
    _text(s, Inches(7.0), Inches(3.9), Inches(5.6), Inches(0.5),
          [("How they're built", 17, GOLD, True)])
    _bullets(s, Inches(7.0), Inches(4.4), Inches(5.6), [
        "Faker (en_IN) + NumPy vectorised generation — seeded & reproducible.",
        "Calibrated vulnerability model (elevation, age, construction, zone).",
        "Vectorised cyclone impact: haversine distance to storm track.",
        "Live enrichment: OpenStreetMap areas/shelters + GDACS cyclones.",
    ], size=13, gap=9, color=WHITE)


def slide_demo(prs):
    s = _slide(prs)
    _kicker(s, "LIVE DEMO FLOW")
    _title(s, "Cyclone enters → agents execute → response allocated")
    steps = [
        ("Cyclone enters region", "Select city + storm (e.g. HUDHUD @ Visakhapatnam); track overlaid on live map."),
        ("Agents execute live", "SSE streams each LangGraph node as it completes — weather → … → judge."),
        ("Risk predicted", "Expected loss, reserve, fraud flags — each with confidence + explainability."),
        ("Resources allocated", "K-Means adjuster zones staged T-24h; personalised customer alerts generated."),
    ]
    y = Inches(2.1)
    for i, (t, d) in enumerate(steps, 1):
        _box(s, Inches(0.7), y, Inches(0.8), Inches(0.8), fill=ACCENT, round_=True)
        _text(s, Inches(0.7), y, Inches(0.8), Inches(0.8), [(str(i), 24, WHITE, True)],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _box(s, Inches(1.7), y, Inches(10.9), Inches(0.95), fill=PANEL,
             line=ACCENT2, line_w=1.0, round_=True)
        _text(s, Inches(1.95), y + Inches(0.12), Inches(10.4), Inches(0.4),
              [(t, 16, WHITE, True)])
        _text(s, Inches(1.95), y + Inches(0.52), Inches(10.4), Inches(0.4),
              [(d, 12, MUTED, False)])
        y = Emu(int(y) + int(Inches(1.15)))


def slide_amd(prs):
    s = _slide(prs)
    _kicker(s, "AMD OPTIMIZATION")
    _title(s, "Tuned for a single MI300X")
    _metric(s, Inches(0.7),  Inches(2.05), "~28 GB", "Qwen3-14B weights", color=ACCENT2)
    _metric(s, Inches(3.55), Inches(2.05), "~135 GB", "KV cache headroom", color=GREEN)
    _metric(s, Inches(6.4),  Inches(2.05), "256", "max concurrent seqs", color=GOLD)
    _metric(s, Inches(9.25), Inches(2.05), "bf16", "native datatype", color=ACCENT)
    _text(s, Inches(0.7), Inches(3.85), Inches(6.0), Inches(0.5),
          [("What we tuned", 17, ACCENT, True)])
    _bullets(s, Inches(0.7), Inches(4.35), Inches(6.0), [
        "gpu-memory-utilization 0.92 + tight max-model-len → more KV slots.",
        "Prefix caching: shared agent system prompt cached once per run.",
        "Chunked prefill → lower p99 under mixed load.",
        "ROCm CK attention kernels (TRITON_FLASH_ATTN=0).",
        "enable_thinking=False → fewer tokens, faster structured output.",
    ], size=13, gap=9)
    _text(s, Inches(7.0), Inches(3.85), Inches(5.6), Inches(0.5),
          [("Throughput (projected, MI300X)", 17, GREEN, True)])
    _bullets(s, Inches(7.0), Inches(4.35), Inches(5.6), [
        "Near-linear scaling to ~128 concurrency.",
        "~230–320 req/s · 4.6–6.4K output tok/s @ 128.",
        "TTFT p50 < 0.15s even under load.",
        "Reproduce: scripts/benchmark_vllm.py → versioned JSON.",
    ], size=13, gap=9, color=WHITE)


def slide_results(prs):
    s = _slide(prs)
    _kicker(s, "RESULTS & IMPACT")
    _title(s, "From reactive claims to proactive catastrophe management")
    cards = [
        ("Faster response", "Adjusters pre-staged in K-Means zones T-24h before landfall — not after.", ACCENT2),
        ("Lower leakage", "FAISS + rule-based fraud flags catch suspicious claims during the surge.", ACCENT),
        ("Reserve accuracy", "Bottom-up IBNR + cat buffer per property, with ±30% wind scenarios.", GOLD),
        ("Regulatory readiness", "LLM-as-Judge audit trail + explainability on every decision.", GREEN),
    ]
    x0, y0 = Inches(0.7), Inches(2.1)
    cw, ch = Inches(5.9), Inches(2.0)
    for i, (t, d, col) in enumerate(cards):
        r, c = divmod(i, 2)
        x = Emu(int(x0) + c * (int(cw) + int(Inches(0.1))))
        y = Emu(int(y0) + r * (int(ch) + int(Inches(0.2))))
        _box(s, x, y, cw, ch, fill=PANEL, line=col, line_w=1.5, round_=True)
        _text(s, x + Inches(0.3), y + Inches(0.25), cw - Inches(0.6), Inches(0.5),
              [(t, 19, col, True)])
        _text(s, x + Inches(0.3), y + Inches(0.85), cw - Inches(0.6), Inches(1.0),
              [(d, 14, WHITE, False)])


def slide_innovation(prs):
    s = _slide(prs)
    _kicker(s, "INNOVATION HIGHLIGHTS")
    _title(s, "What makes IDTCC different")
    items = [
        ("Multi-agent orchestration", "8 specialised agents in a parallel LangGraph DAG — real orchestration, streamed live."),
        ("LLM-as-Judge", "An independent auditor agent scores every forecast on 5 criteria and issues a verdict."),
        ("Deterministic confidence", "Confidence is computed from data, never asked from the model — it can't be hallucinated."),
        ("Guardrails", "JSON validation + numeric bounds + cross-agent consistency catch contradictions."),
        ("Digital-twin simulation", "Per-property risk for 50K twins, aggregated bottom-up."),
        ("Real-time intelligence", "Live OSM + GDACS enrichment; counterfactual track-shift simulator."),
    ]
    x0, y0 = Inches(0.7), Inches(2.0)
    cw, ch = Inches(5.9), Inches(1.35)
    for i, (t, d) in enumerate(items):
        r, c = divmod(i, 2)
        x = Emu(int(x0) + c * (int(cw) + int(Inches(0.1))))
        y = Emu(int(y0) + r * (int(ch) + int(Inches(0.15))))
        col = GOLD if i in (1, 2) else ACCENT2
        _box(s, x, y, cw, ch, fill=PANEL, line=col, line_w=1.25, round_=True)
        _text(s, x + Inches(0.25), y + Inches(0.14), cw - Inches(0.5), Inches(0.4),
              [(t, 15, col, True)])
        _text(s, x + Inches(0.25), y + Inches(0.58), cw - Inches(0.5), Inches(0.7),
              [(d, 12, WHITE, False)])


def slide_roadmap(prs):
    s = _slide(prs)
    _kicker(s, "FUTURE ROADMAP")
    _title(s, "Beyond cyclones")
    phases = [
        ("Now", "Cyclone catastrophe intelligence — 8 agents, 50K twins, MI300X.", ACCENT),
        ("Next", "Floods · earthquakes · wildfires — new hazard engines, same pipeline.", ACCENT2),
        ("Then", "Reinsurance optimization — portfolio-level cession & treaty modelling.", GOLD),
        ("Vision", "Agentic underwriting — real-time, explainable, per-policy pricing.", GREEN),
    ]
    x = Inches(0.7)
    cw = Inches(2.95)
    for name, d, col in phases:
        _box(s, x, Inches(2.4), cw, Inches(3.0), fill=PANEL, line=col, line_w=1.5, round_=True)
        _box(s, x, Inches(2.4), cw, Inches(0.7), fill=col, round_=True)
        _text(s, x, Inches(2.5), cw, Inches(0.5), [(name, 18, BG, True)],
              align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.25), Inches(3.4), cw - Inches(0.5), Inches(1.8),
              [(d, 14, WHITE, False)])
        x = Emu(int(x) + int(cw) + int(Inches(0.1)))


def slide_qr(prs):
    s = _slide(prs)
    _kicker(s, "RESOURCES")
    _title(s, "Explore the project")
    items = [
        ("GitHub", "Full source — backend, frontend, CI"),
        ("Architecture docs", "/docs — architecture, API, AMD, benchmark, ops"),
        ("Demo video", "End-to-end walkthrough"),
    ]
    x = Inches(0.9)
    for t, d in items:
        _box(s, x, Inches(2.3), Inches(3.4), Inches(3.4), fill=PANEL,
             line=ACCENT2, line_w=1.25, round_=True)
        # QR placeholder
        _box(s, x + Inches(0.7), Inches(2.7), Inches(2.0), Inches(2.0),
             fill=WHITE, round_=False)
        _text(s, x + Inches(0.7), Inches(3.5), Inches(2.0), Inches(0.5),
              [("QR", 24, BG, True)], align=PP_ALIGN.CENTER)
        _text(s, x, Inches(4.85), Inches(3.4), Inches(0.5), [(t, 16, WHITE, True)],
              align=PP_ALIGN.CENTER)
        _text(s, x, Inches(5.25), Inches(3.4), Inches(0.5), [(d, 11, MUTED, False)],
              align=PP_ALIGN.CENTER)
        x = Emu(int(x) + int(Inches(3.4)) + int(Inches(0.45)))
    _text(s, Inches(0.9), Inches(6.1), Inches(11), Inches(0.4),
          [("Replace the white squares with generated QR codes before presenting.",
            11, MUTED, False)])


def slide_thanks(prs, team, contact):
    s = _slide(prs)
    _box(s, 0, Inches(3.1), W, Pt(6), fill=ACCENT)
    _text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(1.2),
          [("Thank you", 60, WHITE, True)])
    _text(s, Inches(0.72), Inches(3.4), Inches(12), Inches(0.6),
          [("IDTCC — Insurance Digital Twin Command Center", 20, WHITE, False)])
    _text(s, Inches(0.72), Inches(4.2), Inches(12), Inches(0.5),
          [(f"Team {team}", 16, ACCENT2, True)])
    _text(s, Inches(0.72), Inches(4.75), Inches(12), Inches(0.5),
          [(f"Contact: {contact}", 14, MUTED, False)])
    _chip(s, Inches(0.72), Inches(5.6), "Built on AMD MI300X · ROCm · vLLM",
          w=Inches(5.0), line=ACCENT)


def build(out: str, team: str, contact: str) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    slide_title(prs, team, contact)
    slide_challenge(prs)
    slide_architecture(prs)
    slide_agents(prs)
    slide_twins(prs)
    slide_demo(prs)
    slide_amd(prs)
    slide_results(prs)
    slide_innovation(prs)
    slide_roadmap(prs)
    slide_qr(prs)
    slide_thanks(prs, team, contact)
    # footers (slides 2..11 get numbered; title/thanks kept clean)
    for i, sl in enumerate(prs.slides, 1):
        if 1 < i < 12:
            _footer(sl, i)
    Path(out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"Saved {len(prs.slides)} slides -> {out}")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="IDTCC_Hackathon.pptx")
    ap.add_argument("--team", default="IDTCC")
    ap.add_argument("--contact", default="your-email@example.com")
    args = ap.parse_args()
    build(args.out, args.team, args.contact)
